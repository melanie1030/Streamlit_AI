import streamlit as st
from openai import OpenAI
import pandas as pd
import dotenv  # 用於加載環境變量
import os
from PIL import Image  # 用於圖像處理
import base64  # 用於Base64編碼/解碼
from io import BytesIO  # 用於處理二進制數據
import random  # 用於生成隨機數
import numpy as np
import tempfile
import json
import io
import sys


# 加載環境變量（比如API密鑰）
dotenv.load_dotenv()

# 定義可用的OpenAI模型列表
openai_models = [
    "gpt-4o", 
    "gpt-4-turbo", 
    "gpt-3.5-turbo-16k", 
    "gpt-4", 
    "gpt-4-32k",
]

# 在側邊欄新增資料輸入和保存功能，並加入折疊按鈕
def save_data_as_csv():
    with st.expander("點擊展開以輸入資料並生成 CSV 檔案"):
        input_data = st.text_area("生成 CSV 檔案 (請輸入資料):", height=200)
        
        if input_data:
            try:
                # 將輸入資料處理為表格
                rows = input_data.strip().split("\n")
                data = [row.split(",") for row in rows]

                # 建立資料框架
                df = pd.DataFrame(data[1:], columns=data[0])
                st.write("資料預覽:", df)

                # 提供保存功能
                if st.button("保存為 CSV 檔案"):
                    csv_buffer = io.StringIO()
                    df.to_csv(csv_buffer, index=False)
                    csv_buffer.seek(0)
                    
                    # 下載按鈕
                    st.download_button(
                        label="下載 CSV 檔案",
                        data=csv_buffer.getvalue(),
                        file_name="saved_data.csv",
                        mime="text/csv",
                    )
            except Exception as e:
                st.error(f"資料處理錯誤: {str(e)}")

def execute_code_and_generate_images(code: str, uploaded_file=None):
    try:
        if uploaded_file:
            if uploaded_file.type == "text/csv":
                dataset = pd.read_csv(uploaded_file)  
                globals()['dataset'] = dataset  
            else:
                st.error("只支援CSV檔案")
                return

        import matplotlib.pyplot as plt
        import io

        captured_output = io.StringIO()  
        sys.stdout = captured_output  

        exec(code, globals())  

        sys.stdout = sys.__stdout__

        print_output = captured_output.getvalue()

        image_buffers = []
        figures = [plt.figure(i) for i in plt.get_fignums()]

        for fig in figures:
            buf = io.BytesIO()
            fig.savefig(buf, format="png")  
            buf.seek(0)
            image_buffers.append(buf)

        plt.close('all')

        return print_output, image_buffers

    except Exception as e:
        st.error(f"程式碼執行錯誤: {str(e)}")  # 顯示錯誤訊息
        return None, None

def stream_llm_response(model_params, model_type="openai", api_key=None, uploaded_file=None):
    response_message = ""
    
    try:
        prompt_instructions = """
        當你編寫程式碼時，請始終使用變數 `dataset` 來處理資料集，並避免列印數據集的數據，欄位名稱可以。
        """
        
        dataset_info = ""
        if "dataset" in globals() and isinstance(globals()['dataset'], pd.DataFrame):
            try:
                columns = ", ".join(globals()['dataset'].columns[:100])  
                dataset_info = f"資料集包含 {globals()['dataset'].shape[0]} 行和 {globals()['dataset'].shape[1]} 列。\n欄位名稱：{columns}"
            except Exception as e:
                st.error(f"數據轉換錯誤: {e}")

        if dataset_info:
            st.session_state.messages.append(
                {"role": "user", "content": f"以下是上傳的數據：\n{dataset_info}"}
            )
        
        client = OpenAI(api_key=api_key)
        for chunk in client.chat.completions.create(
            model=model_params["model"],
            messages=[{"role": "system", "content": prompt_instructions}] + st.session_state.messages,
            temperature=model_params["temperature"],
            max_tokens=4096,
            stream=True,
        ):
            chunk_text = chunk.choices[0].delta.content or ""
            response_message += chunk_text
            yield chunk_text

        if "```python" in response_message:
            code_start = response_message.find("```python") + len("```python")
            code_end = response_message.find("```", code_start)
            code = response_message[code_start:code_end].strip()

            print_output, image_buffers = execute_code_and_generate_images(code, uploaded_file)

            if print_output:
                if "程式碼執行錯誤" in print_output:
                    st.error(print_output)
                else:
                    with st.expander("程式碼輸出 (點擊展開)"):
                        st.text_area("程式碼輸出:", value=print_output, height=200)
                
            if isinstance(image_buffers, list):
                for idx, buf in enumerate(image_buffers):
                    with st.expander(f"圖表 {idx + 1} (點擊檢視)"):
                        st.image(buf, caption=f"生成的圖表 {idx + 1}", use_container_width=True)
            else:
                st.error(f"執行程式碼時出現錯誤：{image_buffers}")

    except Exception as e:
        st.error(f"回應處理錯誤: {str(e)}")  # 顯示錯誤訊息


# 將圖像轉換為Base64編碼的函數
def get_image_base64(image_raw):
    buffered = BytesIO()
    image_raw.save(buffered, format=image_raw.format)
    img_byte = buffered.getvalue()
    return base64.b64encode(img_byte).decode('utf-8')

# 將文件轉換為Base64編碼的函數
def file_to_base64(file):
    with open(file, "rb") as f:
        return base64.b64encode(f.read())

# 將Base64編碼轉換回圖像的函數
def base64_to_image(base64_string):
    base64_string = base64_string.split(",")[1]
    return Image.open(BytesIO(base64.b64decode(base64_string)))

# 主函數
def main():
    # --- 頁面配置 ---
    st.set_page_config(
        page_title="Chatbot",
        page_icon="🤖",
        layout="centered",
        initial_sidebar_state="expanded",
    )

    # --- 頁面標題 ---
    st.html("""<h1 style="text-align: center; color: #6ca395;"><i>Openai Chatbot</i> </h1>""")

    # --- 側邊欄設置 ---
    with st.sidebar:
        cols_keys = st.columns(2)
        with cols_keys[0]:
            # 獲取預設的OpenAI API密鑰（如果有的話）
            default_openai_api_key = os.getenv("OPENAI_API_KEY") if os.getenv("OPENAI_API_KEY") is not None else ""
            with st.popover("🔐 API_KEY"):
                openai_api_key = st.text_input("Introduce your OpenAI API Key (https://platform.openai.com/)", 
                                             value=default_openai_api_key, 
                                             type="password")

    # --- 主要內容 ---
    # 檢查用戶是否輸入了API密鑰
    if (openai_api_key == "" or openai_api_key is None or "sk-" not in openai_api_key):
        st.write("#")
        st.warning("⬅️ Please introduce an API Key to continue...")
    
    else:
        client = OpenAI(api_key=openai_api_key)  # 創建OpenAI客戶端

        # 初始化聊天記錄
        if "messages" not in st.session_state:
            st.session_state.messages = []  # 初始化空列表儲存歷史訊息

        # 顯示之前的聊天記錄，並確保每條訊息只渲染一次
        if "rendered_messages" not in st.session_state:
            st.session_state.rendered_messages = set()  # 紀錄已渲染的訊息

        # 顯示聊天記錄
        for idx, message in enumerate(st.session_state.messages):
            # 顯示用戶訊息
            if message["role"] == "user":
                with st.chat_message("user"):
                    # 確保 message["content"] 是列表並且包含字典
                    if isinstance(message["content"], list) and len(message["content"]) > 0:
                        st.markdown(message["content"][0].get("text", ""))  # 使用 get 避免 KeyError
            # 顯示 AI 回應
            elif message["role"] == "assistant":
                with st.chat_message("assistant"):
                    # 確保 message["content"] 是列表並且包含字典
                    if isinstance(message["content"], list) and len(message["content"]) > 0:
                        st.markdown(message["content"][0].get("text", ""))  # 使用 get 避免 KeyError
                        
            # 將已渲染的訊息索引添加到紀錄            
            st.session_state.rendered_messages.add(idx)

        # --- 側邊欄模型選項和輸入 ---
        with st.sidebar:
            st.divider()
            
            # 顯示可用的模型列表
            available_models = [] + (openai_models if openai_api_key else [])
            model = st.selectbox("Select a model:", available_models, index=0)
            model_type = None
            if model.startswith("gpt"): model_type = "openai"
            
            # 模型參數設置
            with st.popover("⚙️ Model parameters"):
                model_temp = st.slider("Temperature", min_value=0.0, max_value=2.0, value=0.3, step=0.1)

            # 設置模型參數
            model_params = {
                "model": model,
                "temperature": model_temp,
            }

            # 重置對話功能
            def reset_conversation():
                if "messages" in st.session_state and len(st.session_state.messages) > 0:
                    st.session_state.pop("messages", None)

            st.button("🗑️", on_click=reset_conversation)

            st.divider()
            
            # 讀取文件內容並返回為文本
            def read_file_content(file):
                if file.type == "application/pdf":
                    from PyPDF2 import PdfReader
                    reader = PdfReader(file)
                    text = ""
                    for page in reader.pages:
                        text += page.extract_text()
                    return text

                elif file.type == "text/plain":
                    return file.getvalue().decode("utf-8")

                elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                    from docx import Document
                    doc = Document(file)
                    text = "\n".join([para.text for para in doc.paragraphs])
                    return text

                elif file.type == "text/csv":
                    # 使用 Pandas 處理 CSV 文件
                    try:
                        file_content = pd.read_csv(file)
                        st.write("### 數據預覽:")
                        st.dataframe(file_content)  # 顯示數據框
                        # 將資料集設定為全域變數，供後續程式碼使用
                        globals()['dataset'] = file_content
                        st.success("CSV 檔案已成功上傳！")
                        return "CSV 文件已載入，請在程式碼中使用 `dataset` 變數進行操作。"
                    except Exception as e:
                        st.error(f"CSV 文件處理出錯: {e}")
                        return None

                elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                    from pptx import Presentation
                    prs = Presentation(file)
                    text = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                    return text
                return None
            
            # 在側邊欄中添加檔案上傳選項
            with st.sidebar:
                st.write("### **Upload a file**:")
                uploaded_file = st.file_uploader(
                    "Upload a file (.txt, .pdf, .docx, .csv, .pptx):", 
                    type=["txt", "pdf", "docx", "csv", "pptx"],
                    accept_multiple_files=False,
                    key="uploaded_file"
                )

                if uploaded_file:
                    file_content = read_file_content(uploaded_file)
                    if file_content:
                        # 將文件內容添加到聊天記錄
                        st.session_state.messages.append(
                            {
                                "role": "user", 
                                "content": [{
                                    "type": "text",  
                                    "text": file_content,
                                }]
                            }
                        )
                    model2key = {
                        "openai": openai_api_key,
                    }
                    image_buffers = stream_llm_response(model_params=model_params,
                                            model_type=model_type,
                                            api_key=model2key[model_type],
                                            uploaded_file=uploaded_file)
            
            # 放置在側邊欄中
            with st.sidebar:
                save_data_as_csv()
                
        # --- 聊天輸入處理 ---
        if prompt := st.chat_input("Hi! Ask me anything..."):
            # 將用戶輸入加入聊天記錄
            st.session_state.messages.append(
                {"role": "user", "content": [{"type": "text", "text": prompt}]}
            )
            # 渲染用戶訊息
            with st.chat_message("user"):
                st.markdown(prompt)

            # 呼叫 GPT 模型，生成回應
            with st.chat_message("assistant"):
                model2key = {
                    "openai": openai_api_key,
                }
                response_text = ""
                assistant_placeholder = st.empty()
                for chunk in stream_llm_response(
                    model_params=model_params,
                    model_type=model_type,
                    api_key=model2key[model_type]
                ):
                    response_text += chunk
                    assistant_placeholder.markdown(response_text)

                # 將 AI 回應加入聊天記錄
                st.session_state.messages.append(
                    {"role": "assistant", "content": [{"type": "text", "text": response_text}]}
                )


# 程式入口點
if __name__=="__main__":
    main()
